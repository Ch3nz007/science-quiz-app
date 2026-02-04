import streamlit as st
import random
import json
import time
from github import Github, GithubException
import google.generativeai as genai
from pypdf import PdfReader
from pptx import Presentation

# --- Setup ---
st.set_page_config(page_title="Quiz Engine", layout="wide")

# ⚠️ CHANGE THIS TO YOUR EXACT REPO NAME
REPO_KEY = "Ch3nz007/science-quiz-app"
FILE_PATH = "science_topics.json"

# --- Cloud Backend (GitHub) ---
def get_repo():
    """Connects to your GitHub Repository"""
    try:
        if "GITHUB_TOKEN" not in st.secrets: return None
        token = st.secrets["GITHUB_TOKEN"]
        g = Github(token)
        return g.get_repo(REPO_KEY)
    except Exception as e:
        return None

def load_data():
    """Reads the JSON file directly from GitHub"""
    try:
        repo = get_repo()
        if not repo: return {"Biology": {}, "Chemistry": {}, "Physics": {}}
        contents = repo.get_contents(FILE_PATH)
        json_content = contents.decoded_content.decode()
        return json.loads(json_content)
    except:
        return {"Biology": {}, "Chemistry": {}, "Physics": {}}

def save_data(data):
    """Updates the JSON file on GitHub"""
    try:
        repo = get_repo()
        if not repo: return
        json_str = json.dumps(data, indent=2)
        try:
            contents = repo.get_contents(FILE_PATH)
            repo.update_file(contents.path, "Update Quiz Data", json_str, contents.sha)
        except:
            repo.create_file(FILE_PATH, "Initial Quiz Data", json_str)
    except Exception as e:
        st.error(f"Failed to save to cloud: {e}")

# --- Helper Functions ---
def extract_text(uploaded_file):
    text = ""
    try:
        if uploaded_file.name.endswith('.pdf'):
            reader = PdfReader(uploaded_file)
            for page in reader.pages: text += page.extract_text() + "\n"
        elif uploaded_file.name.endswith('.pptx'):
            prs = Presentation(uploaded_file)
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"): text += shape.text + "\n"
    except: return None
    return text

def get_available_models(api_key):
    if not api_key: return []
    try:
        genai.configure(api_key=api_key)
        models = genai.list_models()
        return [m.name for m in models if 'generateContent' in m.supported_generation_methods]
    except: return []

# --- CRITICAL FIX: DATA CLEANER ---
def clean_quiz_data(raw_questions):
    """Fixes bad capitalization (Question -> question) and missing keys"""
    if not isinstance(raw_questions, list): return []
    
    cleaned = []
    for q in raw_questions:
        if not isinstance(q, dict): continue
        
        # 1. Normalize Keys (Convert 'Question' to 'question', etc.)
        new_q = {}
        for k, v in q.items():
            new_q[k.lower().strip()] = v
            
        # 2. Ensure Essential Keys Exist
        # If 'question' is missing, try to find a key that looks like it
        if 'question' not in new_q:
            # Fallback: Is there a key named 'prompt' or 'q'?
            if 'prompt' in new_q: new_q['question'] = new_q['prompt']
            else: continue # Skip if truly no question text
            
        # Defaults
        if 'type' not in new_q: new_q['type'] = 'blank'
        if 'options' not in new_q: new_q['options'] = []
        if 'answer' not in new_q: new_q['answer'] = ''
        
        cleaned.append(new_q)
    return cleaned

def generate_questions_gemini(text_content, api_key, model_name):
    genai.configure(api_key=api_key)
    model = genai.GenerativeModel(model_name)
    safe_text = text_content.replace("\\", "/") 
    
    prompt = f"""
    You are a rigorous university examiner. 
    GOAL: Create a massive "Deep Learning" question bank.
    
    INSTRUCTIONS:
    1. Read the text and identify EVERY SINGLE fact, definition, and concept.
    2. For EACH fact found, generate THREE (3) distinct variations:
       - Variation A: Multiple Choice (mcq)
       - Variation B: True/False (true_false)
       - Variation C: Fill in the Blank (blank)

    CRITICAL: 
    - You must generate AT LEAST 60 questions. Do not stop early.
    - JSON Keys must be lowercase: "question", "type", "options", "answer".

    FORMAT: Raw JSON only.
    TEXT CONTENT:
    {safe_text} 
    """
    try:
        response = model.generate_content(prompt)
        cleaned_text = response.text.strip().replace("```json", "").replace("```", "").replace("json\n", "")
        raw_json = json.loads(cleaned_text)
        
        # Run the cleaner before returning
        return clean_quiz_data(raw_json)
        
    except Exception as e:
        st.error(f"AI Generation Error: {e}")
        return None

# --- Frontend ---
st.title("Quiz Engine")

if "quiz_data" not in st.session_state: st.session_state.quiz_data = []
if "score" not in st.session_state: st.session_state.score = 0
if "quiz_active" not in st.session_state: st.session_state.quiz_active = False

data = load_data()

# SIDEBAR
with st.sidebar:
    st.header("Settings")
    subject = st.selectbox("Subject", ["Biology", "Chemistry", "Physics"])
    mode = st.radio("Mode", ["Take Quiz", "Add New Topic", "Manage Topics"])
    st.divider()
    api_key = st.text_input("Gemini API Key", type="password")
    
    available_models = []
    if api_key: available_models = get_available_models(api_key)
    if available_models:
        st.success(f"Connected!")
        selected_model = st.selectbox("Model", available_models, index=0)
    else:
        selected_model = "models/gemini-2.0-flash"

# LOGIC
if mode == "Add New Topic":
    st.subheader(f"Add {subject} Material")
    name = st.text_input("Topic Name")
    uploaded_file = st.file_uploader("Upload Notes", type=["pdf", "pptx"])
    text_input = st.text_area("Paste Text", height=150)
    
    if st.button("Generate & Save"):
        full_text = ""
        if uploaded_file: full_text += extract_text(uploaded_file)
        if text_input: full_text += "\n" + text_input
            
        if full_text.strip() and api_key:
            with st.spinner(f"Generating 60+ Questions..."):
                qs = generate_questions_gemini(full_text, api_key, selected_model)
                if qs:
                    if subject not in data: data[subject] = {}
                    data[subject][name] = qs
                    save_data(data)
                    st.success(f"Success! Saved {len(qs)} questions.")
                    time.sleep(2)
                    st.rerun()
        else:
            st.warning("Missing API Key or Content")

elif mode == "Take Quiz":
    if subject not in data: data[subject] = {}
    
    topics = list(data[subject].keys())
    if topics:
        topic = st.selectbox("Topic", topics)
        questions = data[subject][topic]
        
        if len(questions) > 0:
            q_limit = st.slider("Number of Questions", 1, len(questions), min(10, len(questions)))
            
            if st.button("Start New Quiz"):
                # Ensure we have valid data before sampling
                clean_qs = clean_quiz_data(questions)
                if len(clean_qs) == 0:
                    st.error("Error: This topic has corrupted data. Please delete and regenerate it.")
                else:
                    st.session_state.quiz_data = random.sample(clean_qs, q_limit)
                    st.session_state.score = 0
                    st.session_state.quiz_active = True
                    st.rerun()

            if st.session_state.quiz_active and st.session_state.quiz_data:
                with st.form("my_quiz_form"):
                    user_answers = {}
                    
                    for i, q in enumerate(st.session_state.quiz_data):
                        st.write(f"**{i+1}. {q['question']}**")
                        widget_key = f"q_{i}"
                        
                        q_type = q.get('type', 'blank')
                        if q_type == 'mcq':
                            user_answers[i] = st.radio("Select:", q.get('options', []), key=widget_key, index=None)
                        elif q_type == 'true_false':
                            user_answers[i] = st.radio("True/False:", ["True", "False"], key=widget_key, index=None)
                        else:
                            user_answers[i] = st.text_input("Answer:", key=widget_key)
                        st.divider()
                    
                    submitted = st.form_submit_button("Submit Quiz")
                    
                    if submitted:
                        score = 0
                        st.write("### 📝 Results:")
                        for i, q in enumerate(st.session_state.quiz_data):
                            u_ans = user_answers.get(i)
                            c_ans = q.get('answer', '')
                            is_correct = False
                            q_type = q.get('type', 'blank')
                            
                            if u_ans:
                                if q_type == 'blank':
                                    if str(u_ans).lower().strip() in str(c_ans).lower(): is_correct = True
                                else:
                                    user_str = str(u_ans).split(")")[0].strip()
                                    target_str = str(c_ans).split(")")[0].strip()
                                    if str(u_ans) == str(c_ans) or user_str == target_str:
                                        is_correct = True
                            
                            if is_correct:
                                score += 1
                                st.success(f"**Correct!**")
                            else:
                                st.error(f"**Incorrect**")
                                st.write(f"Your Answer: {u_ans}")
                                st.write(f"Correct Answer: {c_ans}")
                            st.divider()
                            
                        st.metric("Final Score", f"{score}/{len(st.session_state.quiz_data)}")
                        if st.form_submit_button("Take Another Quiz"):
                            st.session_state.quiz_active = False
                            st.rerun()
        else:
            st.warning("Topic has 0 questions.")
    else:
        st.info("No topics found.")

elif mode == "Manage Topics":
    if subject in data:
        for t in list(data[subject].keys()):
            col1, col2 = st.columns([4,1])
            with col1: st.write(f"**{t}** ({len(data[subject][t])} qs)")
            with col2:
                if st.button("Delete", key=f"del_{t}"):
                    del data[subject][t]
                    save_data(data)
                    st.rerun()
